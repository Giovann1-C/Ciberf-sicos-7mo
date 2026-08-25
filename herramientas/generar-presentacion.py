# -*- coding: utf-8 -*-
"""Presentación del Reto — Celda ciberfísica. Genera .pptx con python-pptx."""
import sys
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

DARK   = C(0x12, 0x18, 0x1C); DARK2 = C(0x1B, 0x24, 0x29)
ORANGE = C(0xD8, 0x49, 0x1A); TEAL  = C(0x13, 0x6A, 0x80)
PURPLE = C(0x5C, 0x4A, 0x8A)
LIGHT  = C(0xFF, 0xFF, 0xFF); TINT  = C(0xEF, 0xF3, 0xF2); TINT2 = C(0xE1, 0xE8, 0xE7)
INK    = C(0x12, 0x18, 0x1C); INK2  = C(0x56, 0x64, 0x6A); MUTED = C(0x8A, 0x97, 0x9C)
CRIT   = C(0xAC, 0x2A, 0x21); OK    = C(0x2C, 0x7A, 0x4E); WARN  = C(0xA8, 0x72, 0x1A)
PALE_O = C(0xFB, 0xE7, 0xDF); PALE_C = C(0xFB, 0xEA, 0xE8)
PALE_K = C(0xE3, 0xF1, 0xE8); PALE_T = C(0xDC, 0xEB, 0xEF); PALE_P = C(0xED, 0xE9, 0xF5)
PALE_D = C(0xAE, 0xBD, 0xC2); PALE_L = C(0xD5, 0xDE, 0xE1)
HF, BF, MF = "Arial", "Calibri", "Courier New"

prs = Presentation()
prs.slide_width, prs.slide_height = I(13.333), I(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def txt(s, x, y, w, h, runs, size=14, color=INK, font=BF, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, bullet=False):
    tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        if bullet:
            p.space_after = Pt(8)
        r = p.add_run(); r.text = line
        f = r.font
        f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
        f.color.rgb = color
    return tb


def box(s, x, y, w, h, fill=TINT, line=TINT2, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, I(x), I(y), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = 0.06
    sh.text_frame.text = ""
    return sh


def pill(s, x, y, w, h, label, fill, size=11, color=LIGHT):
    box(s, x, y, w, h, fill=fill, line=None)
    txt(s, x, y, w, h, label, size=size, color=color, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def circle(s, x, y, d, label, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    sh.shadow.inherit = False
    txt(s, x, y, d, d, str(label), size=14, color=LIGHT, bold=True, font=HF,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def arrow(s, x, y, w, h=0.18, fill=MUTED, shape=MSO_SHAPE.RIGHT_ARROW):
    sh = s.shapes.add_shape(shape, I(x), I(y), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    sh.shadow.inherit = False


def head(s, title, sub=None):
    txt(s, 0.75, 0.5, 11.8, 0.75, title, size=32, font=HF, bold=True, color=INK)
    if sub:
        txt(s, 0.75, 1.28, 11.5, 0.5, sub, size=15, color=INK2)


# ============================================================ 1 · portada
s = slide(DARK)
txt(s, 0.9, 1.45, 11.5, 0.35, "RETO MR3005C.601   ·   AGO–DIC 2026   ·   PLAN DE DIRECCIÓN v1.0",
    size=13, color=ORANGE, bold=True)
txt(s, 0.9, 1.95, 11.2, 2.0, ["Celda ciberfísica de", "alimentación robotizada"],
    size=44, font=HF, bold=True, color=LIGHT, spacing=52)
txt(s, 0.9, 4.35, 10.5, 1.2,
    ["AGV omnidireccional + brazo colaborativo sobre riel + visión industrial,",
     "coordinados con un CNC Haas.",
     "",
     "16 semanas   ·   36 % de la calificación   ·   0 días de holgura"],
    size=17, color=PALE_D, spacing=26)

# ============================================================ 2 · la restricción
s = slide()
head(s, "Lo que decide este proyecto son las compras",
     "No la programación. El riel lineal y el LiDAR tardan seis semanas en llegar.")
for i, (n, u, d, col) in enumerate([
        ("6", "semanas", "de lead time en el riel lineal y el LiDAR", ORANGE),
        ("6 sep", "fecha límite", "para emitir las órdenes de compra críticas", CRIT),
        ("0", "días", "de holgura en toda la ruta crítica", TEAL)]):
    x = 0.75 + i * 4.0
    box(s, x, 2.05, 3.7, 2.05)
    txt(s, x + 0.3, 2.2, 3.1, 0.95, n, size=48, font=HF, bold=True, color=col)
    txt(s, x + 0.3, 3.12, 3.1, 0.3, u, size=13, bold=True, color=INK)
    txt(s, x + 0.3, 3.45, 3.1, 0.6, d, size=12, color=INK2)
box(s, 0.75, 4.4, 11.75, 2.1, fill=DARK2, line=None)
txt(s, 1.1, 4.6, 11.0, 0.35, "El escenario que hay que evitar", size=16, font=HF,
    bold=True, color=ORANGE)
txt(s, 1.1, 5.02, 11.0, 1.3,
    ["Si la orden sale al terminar el diseño de detalle —lo natural— el material llega el 8 de noviembre.",
     "Quedan tres semanas para integrar, validar y documentar. No alcanza.",
     "Las fases no se comprimen: se recorta el final, que es donde viven las tres evidencias que más pesan."],
    size=15, color=PALE_L, spacing=25)

# ============================================================ 3 · subsistemas
s = slide()
head(s, "Tres subsistemas, una celda",
     "El AGV entrega en una estación de handoff, no en la máquina. Esa división relaja la repetibilidad exigida al SLAM.")
for i, (t, d, m, col) in enumerate([
        ("AGV omnidireccional",
         "Cuatro ruedas mecanum. Navega por SLAM 2D con LiDAR y Nav2 sobre ROS 2. Transporta el pallet hasta el handoff con docking por AprilTag.",
         "Repetibilidad ≤ ±10 mm", ORANGE),
        ("Brazo sobre riel lineal",
         "Séptimo eje externo de 2.5 m. Toma la pieza del pallet y la deposita en el fixture indexado sobre la mesa del Haas.",
         "Trayectorias enseñadas", TEAL),
        ("Visión industrial",
         "Cámara RGB-D que verifica pass/fail que la pieza esté bien asentada antes de habilitar el ciclo de maquinado.",
         "Respuesta < 3 segundos", PURPLE)]):
    x = 0.75 + i * 4.0
    box(s, x, 2.0, 3.7, 3.85)
    circle(s, x + 0.32, 2.28, 0.46, i + 1, col)
    txt(s, x + 0.32, 2.95, 3.1, 0.65, t, size=17, font=HF, bold=True, color=INK)
    txt(s, x + 0.32, 3.68, 3.1, 1.5, d, size=13, color=INK2)
    txt(s, x + 0.32, 5.3, 3.1, 0.35, m, size=12, bold=True, color=col)
txt(s, 0.75, 6.15, 11.75, 0.45,
    "Control PI / PD / PID sobre microcontrolador NXP S32K312     ·     Coordinación con el CNC por MTConnect, SMB y M-códigos",
    size=14, color=INK2, italic=True, align=PP_ALIGN.CENTER)

# ============================================================ 4 · arquitectura
s = slide()
head(s, "Arquitectura de la celda",
     "El material fluye de izquierda a derecha. El supervisor observa el CNC, pero no lo manda.")
box(s, 4.3, 1.95, 4.6, 0.78, fill=TEAL, line=None)
txt(s, 4.3, 2.07, 4.6, 0.32, "Supervisor de celda", size=15, font=HF, bold=True,
    color=LIGHT, align=PP_ALIGN.CENTER)
txt(s, 4.3, 2.4, 4.6, 0.28, "máquina de estados  ·  dashboard OEE", size=11,
    color=PALE_T, align=PP_ALIGN.CENTER)
txt(s, 4.0, 2.92, 5.2, 0.3, "MQTT  ·  MTConnect :8082  ·  SMB  ·  I/O discreta",
    size=11, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
for i, (t, d, acc) in enumerate([
        ("Estación de carga", "pallet de 5 piezas", False),
        ("AGV mecanum", "SLAM · Nav2 · NXP", True),
        ("Handoff", "AprilTag · conos", False),
        ("Brazo + riel", "7.º eje · 2.5 m", True),
        ("CNC Haas", "fixture indexado", False)]):
    x = 0.75 + i * 2.46
    box(s, x, 3.5, 2.1, 0.98, fill=PALE_O if acc else TINT,
        line=ORANGE if acc else TINT2, lw=2 if acc else 1)
    txt(s, x, 3.66, 2.1, 0.34, t, size=12.5, font=HF, bold=True, color=INK,
        align=PP_ALIGN.CENTER)
    txt(s, x, 4.04, 2.1, 0.3, d, size=10, color=INK2, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + 2.14, 3.9, 0.28)
box(s, 10.59, 4.95, 2.1, 0.82, fill=PALE_P, line=PURPLE, lw=2)
txt(s, 10.59, 5.1, 2.1, 0.32, "Visión RGB-D", size=12, font=HF, bold=True, color=INK,
    align=PP_ALIGN.CENTER)
txt(s, 10.59, 5.44, 2.1, 0.28, "pass / fail", size=10, color=INK2, align=PP_ALIGN.CENTER)
arrow(s, 11.55, 4.52, 0.2, 0.4, PURPLE, MSO_SHAPE.UP_ARROW)
box(s, 0.75, 6.1, 9.5, 0.85)
txt(s, 1.05, 6.28, 8.9, 0.6,
    "El veredicto de visión es una guarda obligatoria: el supervisor solo habilita el ciclo del CNC cuando la cámara aprueba.",
    size=13.5, color=INK)

# ============================================================ 5 · protocolos
s = slide()
head(s, "Por qué tres protocolos y no uno",
     "MTConnect es de solo lectura por diseño. Lo que no puede hacer lo resuelve el canal discreto.")
for i, (cap, proto, dirn, desc, col) in enumerate([
        ("Observación", "MTConnect :8082", "Solo lectura",
         "Estado del CNC cada 250 ms: ejecución, modo, programa activo y alarmas", TEAL),
        ("Transferencia", "Net Share (SMB)", "Escritura",
         "Deposita los programas .nc con hash de contenido para no sobrescribir", PURPLE),
        ("Handshake", "M-códigos + I/O", "Bidireccional",
         "Habilita el ciclo y avisa cuándo terminó. Es lo que MTConnect no puede hacer", ORANGE)]):
    y = 2.05 + i * 1.5
    box(s, 0.75, y, 11.75, 1.3)
    circle(s, 1.05, y + 0.41, 0.48, i + 1, col)
    txt(s, 1.8, y + 0.24, 2.5, 0.35, cap, size=16, font=HF, bold=True, color=INK)
    txt(s, 1.8, y + 0.63, 2.5, 0.3, proto, size=11.5, font=MF, bold=True, color=col)
    txt(s, 4.5, y, 1.9, 1.3, dirn, size=13, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.5, y, 5.7, 1.3, desc, size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.75, 6.75, 11.75, 0.4,
    "La arquitectura es transferible al ecosistema FANUC vía FOCAS 2 sin cambio estructural — argumento de transferencia industrial.",
    size=13, color=INK2, italic=True)

# ============================================================ 6 · cómputo
s = slide()
head(s, "Tres capas de cómputo, un canal de seguridad aparte",
     "Si el software falla, el vehículo se detiene igual. Esa es la única garantía que cuenta.")
for i, (t, d, tag, acc) in enumerate([
        ("Supervisor de celda", "Máquina de estados · MTConnect · SMB · handshake I/O", "orquestación", False),
        ("SBC Linux · ROS 2 Humble", "SLAM · Nav2 · docking AprilTag · cliente MQTT", "tiempo real suave", False),
        ("NXP S32K312", "PWM · encoders · IMU · lazos PI/PD/PID · watchdog", "≥ 200 Hz", True)]):
    y = 2.0 + i * 1.35
    box(s, 0.75, y, 8.0, 1.15, fill=PALE_O if acc else TINT,
        line=ORANGE if acc else TINT2, lw=2 if acc else 1)
    txt(s, 1.05, y + 0.17, 5.6, 0.38, t, size=16, font=HF, bold=True, color=INK)
    txt(s, 1.05, y + 0.6, 6.6, 0.36, d, size=13, color=INK2)
    txt(s, 6.0, y + 0.17, 2.45, 0.35, tag, size=11, bold=True,
        color=ORANGE if acc else TEAL, align=PP_ALIGN.RIGHT)
    if i < 2:
        arrow(s, 4.65, y + 1.17, 0.2, 0.16, MUTED, MSO_SHAPE.DOWN_ARROW)
box(s, 9.15, 2.0, 3.35, 4.05, fill=PALE_C, line=CRIT, lw=2)
txt(s, 9.15, 2.25, 3.35, 0.35, "Canal de seguridad", size=16, font=HF, bold=True,
    color=CRIT, align=PP_ALIGN.CENTER)
txt(s, 9.15, 2.62, 3.35, 0.3, "cableado  ·  independiente", size=11, color=CRIT,
    align=PP_ALIGN.CENTER)
for i, t in enumerate(["Paro de emergencia", "Etapa de potencia"]):
    box(s, 9.5, 3.12 + i * 1.15, 2.65, 0.6, fill=LIGHT, line=CRIT, lw=1.2)
    txt(s, 9.5, 3.12 + i * 1.15, 2.65, 0.6, t, size=13, color=INK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 10.73, 3.78, 0.2, 0.42, CRIT, MSO_SHAPE.DOWN_ARROW)
txt(s, 9.5, 5.15, 2.65, 0.75, ["Corta sin pasar por", "ningún software"], size=12,
    bold=True, color=CRIT, align=PP_ALIGN.CENTER)
txt(s, 0.75, 6.25, 8.0, 0.85,
    "Riesgo R5: cuatro lazos a 200 Hz sobre un solo S32K312. Se mide en la semana 8, no en la 14 — el plan B (segundo MCU) necesita tiempo de compra.",
    size=13, color=INK2)

# ============================================================ 7 · control
s = slide()
head(s, "Cinemática mecanum y lazos de control",
     "Cuatro ruedas a 45°. La velocidad del cuerpo se reparte entre las cuatro.")
box(s, 0.75, 2.0, 6.2, 2.3, fill=DARK2, line=None)
txt(s, 1.05, 2.24, 5.8, 1.85,
    ["w1 = (1/R) · ( vx − vy − (lx+ly)·wz )",
     "w2 = (1/R) · ( vx + vy + (lx+ly)·wz )",
     "w3 = (1/R) · ( vx + vy − (lx+ly)·wz )",
     "w4 = (1/R) · ( vx − vy + (lx+ly)·wz )"],
    size=14, font=MF, color=C(0x9F, 0xD8, 0xE6), spacing=27)
box(s, 7.15, 2.0, 5.35, 2.3, fill=PALE_C, line=CRIT, lw=2)
txt(s, 7.45, 2.2, 4.8, 0.35, "Verifiquen los signos antes de energizar", size=14.5,
    font=HF, bold=True, color=CRIT)
txt(s, 7.45, 2.62, 4.8, 1.55,
    ["Un signo invertido hace que el AGV se mueva en diagonal cuando se le pide avanzar, y el síntoma se confunde con mala sintonía.",
     "",
     "Las mecanum patinan por diseño: la odometría acumula error. Fusionar con IMU y corregir con LiDAR no es opcional."],
    size=12.5, color=INK)
for i, (n, d, w) in enumerate([
        ("PI", "Velocidad de rueda",
         "Elimina error permanente. Sin derivativo: el encoder mete ruido que la derivada amplifica"),
        ("PD", "Posición del riel y orientación",
         "Amortigua sin introducir error de estado estacionario"),
        ("PID", "Posición del AGV en docking",
         "Necesita precisión y amortiguamiento al mismo tiempo")]):
    y = 4.6 + i * 0.72
    pill(s, 0.75, y, 0.95, 0.56, n, TEAL, size=14)
    txt(s, 1.9, y, 3.3, 0.56, d, size=13.5, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.3, y, 7.2, 0.56, w, size=12.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.75, 6.85, 11.75, 0.4,
    "Obligatorios: anti-windup, filtro en el derivativo, watchdog a 200 ms y rampa de aceleración. Cada juego de ganancias se documenta en bitácora.",
    size=12.5, color=INK2, italic=True)

# ============================================================ 8 · ciclo
s = slide()
head(s, "El ciclo de operación",
     "Doce estados. La visión no corrige nada: solo vota, y su voto es obligatorio.")
for i, p in enumerate(["AGV llega al handoff", "Brazo toma la pieza", "Recorre el riel",
                       "Deposita en el fixture", "Retrae a pose segura", "Visión verifica"]):
    x = 0.75 + (i % 3) * 4.0
    y = 2.05 + (i // 3) * 1.35
    v = (i == 5)
    box(s, x, y, 3.7, 1.08, fill=PALE_P if v else TINT, line=PURPLE if v else TINT2,
        lw=2 if v else 1)
    circle(s, x + 0.25, y + 0.32, 0.44, i + 1, PURPLE if v else MUTED)
    txt(s, x + 0.85, y, 2.7, 1.08, p, size=14, bold=v, color=INK, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.75, 4.95, 5.85, 1.15, fill=PALE_K, line=OK, lw=2)
txt(s, 1.05, 5.1, 1.5, 0.35, "PASS", size=15, font=HF, bold=True, color=OK)
txt(s, 1.05, 5.46, 5.3, 0.55,
    "El supervisor habilita el ciclo → el CNC maquina → el brazo retira la pieza → el AGV retorna",
    size=12.5, color=INK)
box(s, 6.85, 4.95, 5.65, 1.15, fill=PALE_C, line=CRIT, lw=2)
txt(s, 7.15, 5.1, 1.5, 0.35, "FAIL", size=15, font=HF, bold=True, color=CRIT)
txt(s, 7.15, 5.46, 5.1, 0.55,
    "Aborta la secuencia, notifica al operador y pide intervención manual. Evita la colisión herramienta–fixture",
    size=12.5, color=INK)
txt(s, 0.75, 6.35, 11.75, 0.4,
    "Cuatro criterios:   presencia   ·   orientación   ·   asentamiento sobre topes   ·   ausencia de objetos extraños",
    size=13.5, color=INK2, align=PP_ALIGN.CENTER)

# ============================================================ 9 · cronograma
s = slide()
head(s, "Dieciséis semanas, cero holgura",
     "Cualquier retraso en compras se traslada íntegro al final.")
X0, ESC = 4.0, 0.685
for i, (cod, nom, fech, off, dur, past, crit) in enumerate([
        ("F0", "Arranque", "10–23 ago", 0.0, 1.1, True, False),
        ("F1", "Definición y compras críticas", "24 ago – 6 sep", 1.1, 1.1, False, True),
        ("F2", "Diseño de detalle", "7–27 sep", 2.2, 1.65, False, False),
        ("F3", "Fabricación", "28 sep – 25 oct", 3.85, 2.2, False, False),
        ("F4", "Navegación", "26 oct – 8 nov", 6.05, 1.1, False, False),
        ("F5", "Celda completa", "9–22 nov", 7.15, 1.1, False, False),
        ("F6", "Cierre", "23 nov – 6 dic", 8.25, 1.05, False, False)]):
    y = 2.1 + i * 0.53
    txt(s, 0.75, y, 3.1, 0.42, cod + "   " + nom, size=12.5, bold=crit,
        color=MUTED if past else INK, anchor=MSO_ANCHOR.MIDDLE)
    box(s, X0 + off * ESC, y + 0.07, max(dur * ESC, 0.3), 0.28,
        fill=ORANGE if crit else (C(0xC9, 0xD2, 0xD5) if past else TEAL), line=None)
    txt(s, 10.55, y, 1.95, 0.42, fech, size=11.5, color=INK2, align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.75, 5.95, 11.75, 1.2, fill=PALE_O, line=ORANGE, lw=2)
txt(s, 1.05, 6.12, 11.2, 0.35, "Hito H1  ·  6 de septiembre  ·  emitir las órdenes de compra críticas",
    size=16, font=HF, bold=True, color=ORANGE)
txt(s, 1.05, 6.52, 11.2, 0.5,
    "Riel lineal, LiDAR, batería, módulo I/O y SBC — $60,500 MXN. Con el diseño congelado solo en interfaces mecánicas y eléctricas, no en el CAD terminado.",
    size=13, color=INK)

# ============================================================ 10 · presupuesto
s = slide()
head(s, "Presupuesto",
     "La meta de la propuesta era menos de $55,000 para el AGV. El AGV va dentro de meta.")
cd = CategoryChartData()
cd.categories = ["Celda (2)", "Visión (1)", "Brazo y riel (2)", "AGV (15 materiales)"]
cd.add_series("MXN", (6500, 7800, 30400, 39340))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, I(0.75), I(2.0), I(7.2), I(3.7), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = False
pl = ch.plots[0]
pl.vary_by_categories = True
pl.has_data_labels = True
dl = pl.data_labels
dl.number_format = '"$"#,##0'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(11); dl.font.color.rgb = INK; dl.font.name = BF
for pt, col in zip(pl.series[0].points, [MUTED, PURPLE, TEAL, ORANGE]):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
ca = ch.category_axis
ca.tick_labels.font.size = Pt(12); ca.tick_labels.font.color.rgb = INK
ca.tick_labels.font.name = BF
ca.has_major_gridlines = False
va = ch.value_axis
va.minimum_scale = 0
va.maximum_scale = 45000
va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = MUTED
va.tick_labels.font.name = BF
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = TINT2
va.major_gridlines.format.line.width = Pt(0.75)
box(s, 8.35, 2.0, 4.15, 3.7)
for i, (n, d, col) in enumerate([("$84,040", "Materiales", INK),
                                 ("$12,600", "Contingencia 15 %", INK2),
                                 ("$96,640", "Total del proyecto", ORANGE)]):
    txt(s, 8.65, 2.3 + i * 1.18, 3.6, 0.6, n, size=29, font=HF, bold=True, color=col)
    txt(s, 8.65, 2.88 + i * 1.18, 3.6, 0.3, d, size=12.5, color=INK2)
txt(s, 0.75, 5.95, 11.75, 0.75,
    "Dónde está el desbordamiento: el riel lineal ($28,000) y la visión ($7,800), que no son parte del AGV. Si hay que recortar, el riel es el candidato — bajar la carrera de 2.5 a 1.5 m y reubicar el handoff.",
    size=14, color=INK)

# ============================================================ 11 · riesgos
s = slide()
head(s, "Los riesgos que pueden costar el semestre",
     "Tres tienen exposición muy alta. Los tres se resuelven esta semana.")
for i, (rid, nom, exp, resp, col) in enumerate([
        ("R7", "Fabricar el brazo consume el semestre", "MUY ALTA",
         "Decidir cobot comercial contra propio antes del 28 de agosto", CRIT),
        ("R6", "El equipo no cierra el BOM a tiempo", "MUY ALTA",
         "Fecha límite dura el 1 de septiembre, revisión diaria en Notion", CRIT),
        ("R1", "El LiDAR se retrasa cuatro semanas", "MUY ALTA",
         "Comprar antes del 6 sep, distribuidor nacional. Nav2 en Gazebo mientras llega", CRIT),
        ("R4", "Sin acceso al CNC Haas en noviembre", "ALTA",
         "Reservar la ventana ya. Plan B: agente MTConnect simulado", WARN),
        ("R2", "El riel no llega o no cumple carrera", "ALTA",
         "Cotizar dos proveedores. Plan B: carrera de 1.5 m y reubicar el handoff", WARN)]):
    y = 2.0 + i * 0.94
    box(s, 0.75, y, 11.75, 0.8)
    txt(s, 1.0, y, 0.7, 0.8, rid, size=15, font=HF, bold=True, color=col,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.75, y, 3.7, 0.8, nom, size=13.5, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    pill(s, 5.6, y + 0.24, 1.2, 0.33, exp, col, size=9)
    txt(s, 7.0, y, 5.3, 0.8, resp, size=12.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.75, 6.8, 11.75, 0.4,
    "Reserva de contingencia: 15 % del presupuesto y una semana de holgura, que solo libera el director de proyecto.",
    size=13, color=INK2, italic=True)

# ============================================================ 12 · decisiones
s = slide()
head(s, "Tres decisiones abiertas",
     "Ninguna es técnica en el fondo. Las tres son de alcance, y encarecen cada semana que pasan sin resolverse.")
for i, (t, f, d, k, col) in enumerate([
        ("¿Brazo propio o cobot comercial?", "28 ago",
         "La propuesta original asume un cobot comercial: UR3e, Doosan o Techman. Fabricar el brazo desde cero es un proyecto completo por sí solo.",
         "Es la decisión más cara del semestre", CRIT),
        ("El NXP no puede correr ROS 2", "1 sep",
         "El S32K312 sirve para los lazos de motor, pero SLAM y Nav2 necesitan Linux. Ya se agregó una SBC al presupuesto.",
         "~$12,000 que no estaban contemplados", WARN),
        ("Mecanum cambia la cinemática", "1 sep",
         "Cuatro motores con encoder controlados independientemente y modelo omnidireccional, no diferencial.",
         "Afecta driver, firmware y presupuesto", WARN)]):
    x = 0.75 + i * 4.0
    box(s, x, 2.05, 3.7, 4.1)
    pill(s, x + 0.3, 2.35, 1.35, 0.4, f, col, size=11)
    txt(s, x + 0.3, 2.92, 3.1, 0.95, t, size=16, font=HF, bold=True, color=INK)
    txt(s, x + 0.3, 3.95, 3.1, 1.6, d, size=12.5, color=INK2)
    txt(s, x + 0.3, 5.55, 3.1, 0.5, k, size=12.5, bold=True, color=col)

# ============================================================ 13 · sección
s = slide(DARK)
circle(s, 0.9, 2.35, 0.95, "2", ORANGE)
txt(s, 2.2, 2.35, 9.6, 0.95, "Cómo vamos a trabajar", size=34, font=HF, bold=True,
    color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 2.2, 3.5, 9.6, 0.9,
    "Notion para el equipo  ·  Teams para la conversación  ·  Canvas para las entregas  ·  Claude como agente de gestión",
    size=16, color=PALE_D)

# ============================================================ 14 · flujo
s = slide()
head(s, "El flujo de trabajo",
     "Nadie comparte una cuenta y ningún dato vive en dos lugares a la vez.")
for i, (t, d, fill, line) in enumerate([
        ("Equipo", "4 integrantes", TINT, TINT2),
        ("Notion", "tareas · BOM · riesgos", PALE_T, TEAL),
        ("Rodrigo", "director de proyecto", TINT, TINT2),
        ("Claude", "agente / project manager", PALE_O, ORANGE)]):
    x = 0.75 + i * 3.1
    box(s, x, 2.3, 2.7, 1.15, fill=fill, line=line, lw=2)
    txt(s, x, 2.52, 2.7, 0.4, t, size=17, font=HF, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, x, 2.95, 2.7, 0.3, d, size=11, color=INK2, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + 2.76, 2.79, 0.28)
        txt(s, x + 2.5, 3.05, 0.8, 0.3, ["captura", "valida", "pide"][i], size=10,
            color=MUTED, align=PP_ALIGN.CENTER)
for i, t in enumerate(["PO  ·  OF  ·  hojas viajeras", "Guías  ·  briefs  ·  análisis"]):
    box(s, 9.85, 3.95 + i * 0.8, 2.65, 0.65, fill=LIGHT, line=ORANGE, lw=1.2)
    txt(s, 9.85, 3.95 + i * 0.8, 2.65, 0.65, t, size=11.5, color=INK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 11.08, 3.5, 0.2, 0.4, ORANGE, MSO_SHAPE.DOWN_ARROW)
box(s, 0.75, 5.8, 11.75, 1.15, fill=DARK2, line=None)
txt(s, 1.05, 6.1, 11.2, 0.6,
    "Claude genera el documento; Rodrigo autoriza. El equipo nunca toca la cuenta y aun así todo lo que captura llega a los entregables.",
    size=15, color=PALE_L)

# ============================================================ 15 · qué hace / no hace
s = slide()
head(s, "Claude como agente y project manager",
     "Lo que sí hace y lo que deliberadamente no hace.")
box(s, 0.75, 2.0, 5.85, 4.3, fill=PALE_K, line=OK, lw=2)
txt(s, 1.1, 2.25, 5.2, 0.4, "Lo que hace", size=19, font=HF, bold=True, color=OK)
txt(s, 1.1, 2.78, 5.2, 3.4,
    ["·  Emite PO, órdenes de fabricación y hojas viajeras en el formato del profesor",
     "·  Calcula fechas de llegada por lead time y avisa qué está en riesgo",
     "·  Brief automático de lunes a sábado a las 7:03",
     "·  Responde dudas citando archivo y página de un índice de 773 páginas",
     "·  Transcribe las libretas de GoodNotes y las conecta con el proyecto",
     "·  Sincroniza Notion, archivos locales y Google Calendar"],
    size=13.5, color=INK, bullet=True)
box(s, 6.9, 2.0, 5.6, 4.3, fill=PALE_C, line=CRIT, lw=2)
txt(s, 7.25, 2.25, 4.9, 0.4, "Lo que no hace", size=19, font=HF, bold=True, color=CRIT)
txt(s, 7.25, 2.78, 4.9, 3.4,
    ["·  No decide el alcance. Señala las divergencias; la decisión es del equipo",
     "·  No compra ni compromete dinero. Genera la orden; la autoriza Rodrigo",
     "·  No sustituye al director de proyecto. Rodrigo aprueba y responde",
     "·  No lo usa el equipo directamente. Una sola cuenta, sin compartir"],
    size=13.5, color=INK, bullet=True)
txt(s, 0.75, 6.55, 11.75, 0.45,
    "Un agente que recuerda, documenta y vigila plazos. No un sustituto del criterio de ingeniería.",
    size=14.5, color=INK2, italic=True, align=PP_ALIGN.CENTER)

# ============================================================ 16 · ritmo
s = slide()
head(s, "El ritmo semanal", "Seis momentos fijos. Todo lo demás es trabajo técnico.")
for i, (c, q, w) in enumerate([
        ("Lunes AM", "Junta de avance y actualización de Notion", "Equipo"),
        ("Lunes AM", "Planeación de la semana y estado del reto", "Rodrigo"),
        ("Diario 7:03", "Brief automático con alertas de plazo", "Claude"),
        ("Al cotizar", "Actualizar proveedor, costo y lead time en el BOM", "Quien cotiza"),
        ("Al aprobar", "Emitir la orden de compra", "Rodrigo"),
        ("Viernes", "Revisión de riesgos y avance contra el plan", "Rodrigo")]):
    y = 2.0 + i * 0.78
    box(s, 0.75, y, 11.75, 0.66)
    txt(s, 1.05, y, 1.9, 0.66, c, size=13, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.1, y, 7.2, 0.66, q, size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 10.3, y, 1.95, 0.66, w, size=12.5, color=INK2, align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 17 · cierre
s = slide(DARK)
txt(s, 0.9, 1.2, 11.5, 0.8, "Lo que pasa esta semana", size=38, font=HF, bold=True, color=LIGHT)
for i, (f, a, col) in enumerate([
        ("28 ago", "Decidir brazo propio contra cobot comercial", CRIT),
        ("31 ago", "WBS completo y matriz RACI del equipo", WARN),
        ("1 sep", "Cerrar el BOM preliminar — bloquea todas las compras", CRIT),
        ("2 sep", "Cotizar riel y LiDAR con dos proveedores cada uno", WARN),
        ("6 sep", "Emitir las cinco órdenes de compra críticas", ORANGE),
        ("6 sep", "Reservar la ventana de acceso al CNC Haas", WARN)]):
    y = 2.3 + i * 0.72
    pill(s, 0.9, y, 1.4, 0.52, f, col, size=13)
    txt(s, 2.55, y, 9.8, 0.52, a, size=16, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.9, 6.75, 11.5, 0.5,
    "Si el laboratorio no da acceso al Haas en noviembre, no hay celda que demostrar. Eso no se arregla con dinero ni con horas.",
    size=14, italic=True, color=PALE_D)

prs.save(sys.argv[1])
print("OK  %d diapositivas  ->  %s" % (len(prs.slides._sldIdLst), sys.argv[1]))
